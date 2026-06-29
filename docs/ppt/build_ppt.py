# -*- coding: utf-8 -*-
"""
PowerPMAC MCP 교육자료(.pptx) 생성 스크립트.

구조와 동작 원리를 가르치는 종합 교육 슬라이드를 python-pptx로 코드 생성한다.
소스(이 파일)가 곧 슬라이드 정의 — 수정 후 재실행하면 .pptx가 재생성된다.

  pip install python-pptx
  python build_ppt.py   ->  PowerPMAC_MCP_교육자료.pptx

출처 매핑:
  개요/사용법/워크플로  : docs/manual.md, README.md
  Skill                : Skills/powerpmac-dev/SKILL.md
  MCP 프로토콜          : mcp-server/Program.cs
  MCP 브리지            : mcp-server/PmacBridge.cs, mcp-server/README.md
  PDK 런타임/이식성     : mcp-server/PdkRuntime.cs, PowerPmacMcp.csproj
  설치/등록             : setup.ps1, INSTALL.md
"""
import os
import re
import sys

from pptx import Presentation
from pptx.util import Inches as In, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

# ---------------------------------------------------------------- 테마/상수
PRIMARY = RGBColor(0x16, 0x32, 0x4F)   # 네이비 (제목/주색)
PRIMARY2 = RGBColor(0x21, 0x4E, 0x79)  # 밝은 네이비
ACCENT = RGBColor(0xE8, 0x6A, 0x17)    # 주황 (강조)
SKILL_C = RGBColor(0x2E, 0x7D, 0x32)   # 초록 (Skill)
MCP_C = RGBColor(0x15, 0x65, 0xC0)     # 파랑 (MCP)
TEAL_C = RGBColor(0x00, 0x83, 0x8F)    # 청록 (브리지)
RT_C = RGBColor(0x6A, 0x3D, 0x9A)      # 보라 (런타임)
CTRL_C = RGBColor(0x4D, 0x54, 0x5B)    # 회색 (컨트롤러)
INK = RGBColor(0x24, 0x29, 0x2E)       # 본문 잉크
GRAY_D = RGBColor(0x5A, 0x63, 0x6B)    # 진회색
PANEL = RGBColor(0xF1, 0xF5, 0xF9)     # 패널 배경
PANEL2 = RGBColor(0xE3, 0xEA, 0xF1)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CODE_BG = RGBColor(0x1E, 0x24, 0x2B)   # 코드 배경
CODE_FG = RGBColor(0xDC, 0xE3, 0xEA)   # 코드 글자
RED_C = RGBColor(0xC0, 0x39, 0x2B)     # 경고(빨강)

FONT_KR = "맑은 고딕"
FONT_CODE = "Consolas"

prs = Presentation()
prs.slide_width = In(13.333)
prs.slide_height = In(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height
MARGIN = In(0.6)
CW = SW - In(1.2)
PAGE = [0]
FOOTER_TXT = "Power PMAC 개발 도구 · 구조와 동작 원리"

# ---------------------------------------------------------------- 저수준 헬퍼
def set_font(run, name=FONT_KR, ea=None):
    """라틴+동아시아 폰트를 함께 지정해 한글 깨짐 방지.
    ea를 따로 주면 동아시아(한글) 글리프만 다른 폰트로 — 코드박스에서
    영문은 Consolas, 한글 주석은 맑은 고딕으로 분리 렌더하는 데 쓴다."""
    ea_name = ea or name
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    latin = rPr.find(qn('a:latin'))
    if latin is None:
        latin = rPr.makeelement(qn('a:latin'), {})
        rPr.insert(0, latin)
    latin.set('typeface', name)
    after = latin
    for tag, val in (('a:ea', ea_name), ('a:cs', name)):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            after.addnext(el)
        el.set('typeface', val)
        after = el


def add_rich(p, text, size, color, font=FONT_KR, bold_default=False, accent=ACCENT):
    """인라인 마크업 파서:
      **...**  → 강조(주황+볼드)
      `...`    → 코드(고정폭 Consolas, 문맥 색 유지) — 백슬래시 경로도 정상 표시
    """
    for part in re.split(r'(\*\*.+?\*\*|`[^`]+`)', str(text)):
        if not part:
            continue
        if len(part) >= 4 and part.startswith('**') and part.endswith('**'):
            seg, emph, code = part[2:-2], True, False
        elif len(part) >= 2 and part.startswith('`') and part.endswith('`'):
            seg, emph, code = part[1:-1], False, True
        else:
            seg, emph, code = part, False, False
        if not seg:
            continue
        r = p.add_run()
        r.text = seg
        r.font.size = Pt(size)
        r.font.bold = True if emph else bold_default
        r.font.color.rgb = accent if emph else color
        set_font(r, FONT_CODE if code else font, ea=FONT_KR if code else None)


def bg(slide, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def textbox(slide, x, y, w, h, lines, size=18, color=INK, bold=False,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT_KR, gap=4, line_spacing=1.05):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(lines, str):
        lines = [lines]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(gap)
        p.line_spacing = line_spacing
        add_rich(p, ln, size, color, font, bold_default=bold)
    return tb


def title_bar(slide, title, kicker=None):
    """상단 제목 + 강조 밑줄."""
    if kicker:
        textbox(slide, MARGIN, In(0.34), CW, In(0.32), kicker, size=13,
                color=ACCENT, bold=True)
        ty = In(0.62)
    else:
        ty = In(0.42)
    textbox(slide, MARGIN, ty, CW, In(0.7), title, size=28, color=PRIMARY, bold=True)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, In(1.32), In(1.0), In(0.07))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    bar.shadow.inherit = False


def footer(slide):
    textbox(slide, MARGIN, In(7.06), In(9.0), In(0.3), FOOTER_TXT, size=9, color=GRAY_D)
    textbox(slide, SW - In(1.4), In(7.06), In(0.8), In(0.3), str(PAGE[0]),
            size=9, color=GRAY_D, align=PP_ALIGN.RIGHT)


def content(title, kicker=None):
    slide = prs.slides.add_slide(BLANK)
    PAGE[0] += 1
    title_bar(slide, title, kicker)
    footer(slide)
    return slide


def bullets(slide, items, x=MARGIN, y=In(1.6), w=CW, h=In(5.2), size=18, gap=8):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        lvl, text = item if isinstance(item, tuple) else (0, item)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap if lvl == 0 else gap - 3)
        p.line_spacing = 1.16
        pre = p.add_run()
        pre.text = "•  " if lvl == 0 else "       –  "
        pre.font.size = Pt(size if lvl == 0 else size - 2)
        pre.font.bold = (lvl == 0)
        pre.font.color.rgb = PRIMARY2 if lvl == 0 else GRAY_D
        set_font(pre)
        add_rich(p, text, size if lvl == 0 else size - 2, INK if lvl == 0 else GRAY_D)
    return tb


def table(slide, data, x=MARGIN, y=In(1.7), w=CW, h=None, col_w=None,
          size=13, header=True, row_h=In(0.42)):
    rows, cols = len(data), len(data[0])
    if h is None:
        h = row_h * rows
    gt = slide.shapes.add_table(rows, cols, x, y, w, h).table
    gt.first_row = header
    gt.horz_banding = False
    if col_w:
        total = sum(col_w)
        for c in range(cols):
            gt.columns[c].width = Emu(int(int(w) * col_w[c] / total))
    for r in range(rows):
        gt.rows[r].height = row_h
        for c in range(cols):
            cell = gt.cell(r, c)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = In(0.1)
            cell.margin_right = In(0.06)
            cell.margin_top = In(0.02)
            cell.margin_bottom = In(0.02)
            cell.fill.solid()
            if header and r == 0:
                cell.fill.fore_color.rgb = PRIMARY
            else:
                cell.fill.fore_color.rgb = WHITE if (r % 2 == 1) else PANEL
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            col = WHITE if (header and r == 0) else INK
            add_rich(p, data[r][c], size, col, bold_default=(header and r == 0))
    return gt


def code_box(slide, code_lines, x=MARGIN, y=In(1.7), w=CW, h=In(2.0), size=12.5):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = CODE_BG
    shp.line.fill.background()
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.word_wrap = False          # 코드는 줄바꿈하지 않는다(라인 정렬 보존)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = In(0.18)
    tf.margin_right = In(0.12)
    tf.margin_top = In(0.12)
    tf.margin_bottom = In(0.1)
    for i, ln in enumerate(code_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.05
        p.alignment = PP_ALIGN.LEFT   # 도형 첫 문단의 기본 중앙정렬 방지
        r = p.add_run()
        r.text = ln if ln != "" else " "
        r.font.size = Pt(size)
        r.font.color.rgb = CODE_FG
        set_font(r, FONT_CODE, ea=FONT_KR)
    return shp


# ---------------------------------------------------------------- 다이어그램 헬퍼
def box(slide, x, y, w, h, lines, fill, txt=WHITE, size=13, bold=True,
        align=PP_ALIGN.CENTER, rounded=True, line=None, font=FONT_KR):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is not None:
        shape.line.color.rgb = line
        shape.line.width = Pt(1.25)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = In(0.06)
    tf.margin_right = In(0.06)
    tf.margin_top = In(0.03)
    tf.margin_bottom = In(0.03)
    if isinstance(lines, str):
        lines = [lines]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = 1.0
        add_rich(p, ln, size, txt, font, bold_default=bold)
    return shape


def diamond(slide, x, y, w, h, lines, fill=ACCENT, txt=WHITE, size=12):
    shape = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    if isinstance(lines, str):
        lines = [lines]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.line_spacing = 1.0
        add_rich(p, ln, size, txt, FONT_KR, bold_default=True)
    return shape


def arrow(slide, x1, y1, x2, y2, color=GRAY_D, width=2.0, dashed=False):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    ln = conn.line._get_or_add_ln()
    if dashed:
        dash = ln.makeelement(qn('a:prstDash'), {'val': 'dash'})
        ln.append(dash)
    tail = ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
    ln.append(tail)
    return conn


def line(slide, x1, y1, x2, y2, color=GRAY_D, width=1.5, dashed=False):
    """화살표 머리 없는 직선 — 다단(ㄷ자) 경로의 중간 구간용."""
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    if dashed:
        ln = conn.line._get_or_add_ln()
        ln.append(ln.makeelement(qn('a:prstDash'), {'val': 'dash'}))
    return conn


def alabel(slide, x, y, w, text, color=GRAY_D, size=11, align=PP_ALIGN.CENTER):
    textbox(slide, x, y, w, In(0.3), text, size=size, color=color, align=align, bold=True)


# ================================================================ 슬라이드 빌드
def s_cover():
    slide = prs.slides.add_slide(BLANK)
    bg(slide, PRIMARY)
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, In(4.55), SW, In(0.09))
    band.fill.solid(); band.fill.fore_color.rgb = ACCENT
    band.line.fill.background(); band.shadow.inherit = False
    textbox(slide, In(1.0), In(2.0), In(11.3), In(0.5),
            "OMRON Delta Tau Power PMAC", size=18, color=RGBColor(0xBF, 0xD3, 0xE6), bold=True)
    textbox(slide, In(1.0), In(2.5), In(11.3), In(1.3),
            "개발 도구: 구조와 동작 원리", size=46, color=WHITE, bold=True)
    textbox(slide, In(1.0), In(3.8), In(11.3), In(0.6),
            "Claude Code Skill + MCP 서버", size=24, color=RGBColor(0xE8, 0x6A, 0x17), bold=True)
    textbox(slide, In(1.0), In(4.8), In(11.3), In(0.9),
            ["자연어로 Power PMAC 코드를 작성하고, 실제 컨트롤러를 빌드·다운로드·제어한다.",
             "powerpmac-dev (Skill, 지식)  +  powerpmac (MCP, 조작)"],
            size=15, color=RGBColor(0xCF, 0xDB, 0xE8))
    textbox(slide, In(1.0), In(6.5), In(11.3), In(0.5),
            "2026-06  ·  교육자료", size=13, color=RGBColor(0x9F, 0xB3, 0xC8))


def s_toc():
    slide = content("목차")
    rows = [
        ["Part 1", "개요 — 무엇을, 왜", "문제의식 · Skill+MCP · 전체 구성도"],
        ["Part 2", "Skill 계층 — Power PMAC 지식", "자동 적용 · 지식 라우팅 · 멘탈모델 · 안전 함정"],
        ["Part 3", "MCP 계층 — 컨트롤러 조작 (동작 원리 핵심)", "3계층 아키텍처 · 프로토콜 · 빌드/다운로드/세션 · 런타임"],
        ["Part 4", "설치 · 배포 · 운영", "setup.ps1 · 듀얼 리모트 · 개발 워크플로 · 실증"],
        ["Part 5", "마무리", "핵심 설계 결정 요약 · 참고자료"],
    ]
    table(slide, rows, y=In(1.8), col_w=[1.1, 4.2, 5.0], size=15, header=False, row_h=In(0.92))


def s_divider(part, title, sub):
    slide = prs.slides.add_slide(BLANK)
    bg(slide, PRIMARY)
    side = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, In(0.28), SH)
    side.fill.solid(); side.fill.fore_color.rgb = ACCENT
    side.line.fill.background(); side.shadow.inherit = False
    textbox(slide, In(1.1), In(2.6), In(11), In(0.7), part, size=22,
            color=ACCENT, bold=True)
    textbox(slide, In(1.1), In(3.2), In(11), In(1.2), title, size=40, color=WHITE, bold=True)
    textbox(slide, In(1.1), In(4.5), In(11), In(0.6), sub, size=16,
            color=RGBColor(0xC0, 0xD0, 0xE0))


# ---- Part 1 ----
def s_problem():
    slide = content("문제의식 — 기존 Power PMAC 개발의 마찰", "PART 1 · 개요")
    bullets(slide, [
        "**방대한 매뉴얼** — 요소·명령이 버전에 민감. 공식 매뉴얼 4종 약 **3,260쪽**을 매번 뒤져야 한다.",
        "**수동 개발 루프** — IDE에서 빌드 → 다운로드 → 터미널 조회를 사람이 일일이 반복.",
        "**안전 함정** — kill ≠ abort, 안전장치 기본 OFF, 모터·축 단위 혼동 … 작은 실수가 장비 사고로.",
        "**반복 조회** — 여러 요소 값 확인, 셸 로그 점검 등 손이 많이 가는 작업.",
        ("", "→ AI 에이전트가 **지식**(정확한 코드)과 **조작**(실제 장비)을 모두 도와준다면?"),
    ], size=19, gap=14)


def s_solution():
    slide = content("해결책 — Claude Code + Skill + MCP", "PART 1 · 개요")
    bullets(slide, [
        "**Claude Code** = 터미널·IDE에서 동작하는 AI 코딩 에이전트. 자연어로 시키면 파일을 읽고·쓰고·실행.",
        "여기에 두 가지를 더한다:",
        (1, "**Skill** `powerpmac-dev` — Power PMAC을 '이해'. 문법·데이터구조·함정을 자동 적용해 정확한 코드·리뷰."),
        (1, "**MCP** `powerpmac` — 컨트롤러를 '조작'. 빌드·다운로드·라이브 조회·셸을 자연어로."),
    ], size=18, h=In(2.4), gap=10)
    box(slide, MARGIN, In(4.5), In(5.9), In(2.0),
        ["Before", "IDE 열기 → 수동 빌드 → 수동 다운로드", "→ 매뉴얼 검색 → 터미널로 값 확인",
         "(여러 도구를 오가며 반복)"],
        fill=PANEL, txt=INK, size=14, bold=False, rounded=True, line=PANEL2)
    box(slide, In(6.8), In(4.5), In(5.9), In(2.0),
        ["After", "\"이 프로젝트 빌드해서 192.168.0.200에", "다운로드하고 Motor[1].ActPos 읽어줘\"",
         "— 한 마디로 전체 루프"],
        fill=SKILL_C, txt=WHITE, size=14, bold=False)
    box(slide, In(6.45), In(4.95), In(0.5), In(0.5), "▶", fill=ACCENT, size=16)


def s_two_axes():
    slide = content("두 축 — Skill(지식)과 MCP(조작)", "PART 1 · 개요")
    rows = [
        ["구분", "Skill — powerpmac-dev", "MCP — powerpmac"],
        ["역할", "**이해** · 코드 작성 · 리뷰 · 디버깅", "**조작** · 실제 컨트롤러 제어"],
        ["형태", "마크다운 지식베이스 (~/.claude/skills)", ".NET stdio 서버 (powerpmac-mcp.exe)"],
        ["호출 방식", "질문하면 **자동 적용** (부를 필요 없음)", "\"~해줘\" 하면 Claude가 **도구 호출**"],
        ["필요 환경", "PDK 불필요 (-SkillOnly 로 단독)", "PDK + Power PMAC IDE **필수**"],
        ["예", "\"조그 모션 프로그램 짜줘\"", "\"빌드·다운로드·조그 해줘\""],
    ]
    table(slide, rows, y=In(1.7), col_w=[1.4, 4.8, 4.8], size=14, row_h=In(0.62))
    textbox(slide, MARGIN, In(6.0), CW, In(0.6),
            "→ 둘은 상보적이다: **Skill로 작성하고, MCP로 돌린다.**", size=16, color=PRIMARY, bold=True)


def s_sysdiagram():
    slide = content("전체 시스템 구성도", "PART 1 · 개요")
    # 개발자
    box(slide, In(0.7), In(3.1), In(1.7), In(1.0), ["개발자", "(자연어 지시)"], fill=CTRL_C, size=14)
    arrow(slide, In(2.4), In(3.6), In(3.05), In(3.6), color=PRIMARY2, width=2.5)
    # Claude Code
    box(slide, In(3.05), In(2.95), In(2.7), In(1.3), ["Claude Code", "AI 에이전트", "(MCP 클라이언트)"],
        fill=PRIMARY, size=15)
    # Skill (위)
    box(slide, In(6.5), In(1.7), In(3.1), In(1.25),
        ["Skill: powerpmac-dev", "매뉴얼 ~3,260쪽 정제 지식", "(코드 작성·리뷰)"], fill=SKILL_C, size=13)
    arrow(slide, In(5.75), In(3.3), In(6.5), In(2.4), color=SKILL_C, width=2.5)
    # MCP (아래)
    box(slide, In(6.5), In(3.95), In(3.1), In(1.25),
        ["MCP 서버", "powerpmac-mcp.exe", "(net48 · x86)"], fill=MCP_C, size=13)
    arrow(slide, In(5.75), In(3.9), In(6.5), In(4.5), color=MCP_C, width=2.5)
    # Skill → MCP : 작성한 코드를 빌드·다운로드 (연동 흐름)
    arrow(slide, In(8.05), In(2.95), In(8.05), In(3.9), color=SKILL_C, width=2.0)
    alabel(slide, In(8.2), In(3.28), In(1.7), "작성 코드 →", size=9, color=SKILL_C, align=PP_ALIGN.LEFT)
    # Controller
    box(slide, In(10.25), In(3.95), In(2.4), In(1.25),
        ["Power PMAC", "컨트롤러", "192.168.0.200"], fill=CTRL_C, size=13)
    arrow(slide, In(9.6), In(4.55), In(10.25), In(4.55), color=MCP_C, width=2.5)
    alabel(slide, In(9.5), In(4.2), In(1.0), "SSH/gpascii", size=10)
    # PDK
    box(slide, In(6.5), In(5.7), In(3.1), In(0.85),
        ["PDK DLL · 크로스컴파일러", "(ODT.* · cygwin · rsync)"], fill=RT_C, size=12)
    arrow(slide, In(8.05), In(5.7), In(8.05), In(5.2), color=RT_C, width=2.0)
    alabel(slide, In(9.05), In(5.55), In(1.6), "인프로세스 로드", size=10, align=PP_ALIGN.LEFT)
    textbox(slide, MARGIN, In(6.74), CW, In(0.35),
            "Claude가 Skill·MCP를 자동 선택 — Skill로 작성한 코드를 MCP가 빌드·다운로드해 컨트롤러에 반영한다.",
            size=12, color=GRAY_D)


# ---- Part 2 ----
def s_skill_mech():
    slide = content("Skill 메커니즘 — 질문하면 자동 적용", "PART 2 · Skill")
    bullets(slide, [
        "**Skill** = 특정 도메인의 지식·절차를 담은 Claude Code 모듈. 질문이 도메인에 맞으면 **자동 로드**.",
        "설치 위치: **~/.claude/skills/powerpmac-dev** (setup.ps1 이 연결).",
        "트리거(description): \"Power PMAC / PMAC / Delta Tau / Sysmac 모션, prog·plc·C, 데이터구조 …\" 관련이면 발동.",
        "**점진적 로드로 토큰 절약** — 항상 읽는 건 SKILL.md(맵)뿐. 깊은 내용은 필요 시 reference/만 연다.",
        "사용자는 Skill을 **명시적으로 부르지 않는다** — 그냥 Power PMAC을 물으면 된다.",
    ], size=18, gap=13)


def s_skill_struct():
    slide = content("powerpmac-dev 구조 — 맵 → 정제 → 원문", "PART 2 · Skill")
    box(slide, MARGIN, In(1.7), In(3.9), In(1.5),
        ["SKILL.md", "맵 + 안전 요약", "(멘탈모델 · top gotchas · 라우팅 표)", "항상 읽힘"],
        fill=SKILL_C, size=13)
    arrow(slide, In(4.5), In(2.45), In(5.0), In(2.45), color=GRAY_D)
    box(slide, In(5.0), In(1.7), In(3.9), In(1.5),
        ["reference/*.md", "도메인별 정제 distillation", "syntax · data-structure · motion ·", "plc · c-programming · gotchas …"],
        fill=MCP_C, size=13)
    arrow(slide, In(8.9), In(2.45), In(9.4), In(2.45), color=GRAY_D)
    box(slide, In(9.4), In(1.7), In(3.3), In(1.5),
        ["reference/raw/", "공식 매뉴얼 원문", "~3,260쪽", "(grep 최후 근거)"],
        fill=CTRL_C, size=13)
    textbox(slide, MARGIN, In(3.5), CW, In(0.4), "동작 규율 (정확성 우선):", size=16, color=PRIMARY, bold=True)
    bullets(slide, [
        "요소·명령 이름이 불확실하면 **raw를 grep해 확인** — 추측으로 만들어내지 않는다.",
        "정제 노트에 \"(verify: …)\" 표시가 있으면 원문으로 재확인 후 사용.",
        "비자명한 규칙을 말할 땐 **매뉴얼 페이지를 인용** — 사용자가 검증 가능.",
    ], y=In(3.95), size=17, gap=10)


def s_routing():
    slide = content("지식 라우팅 — 도메인 → reference 파일", "PART 2 · Skill")
    rows = [
        ["필요한 것", "reference 파일"],
        ["연산자 · 변수(P/Q/M/L/I) · 흐름제어 · 온라인 vs 버퍼 명령", "syntax-rules.md"],
        ["Structure[index].Element 모델, Sys./Motor[]/Coord[]/Gate3[], SAVED", "data-structure.md"],
        ["모션: open prog, 무브모드, 축정의, 좌표계, 룩어헤드, G코드", "script-motion.md"],
        ["PLC: open plc, 스캔모델, 타이머, cmd, 시퀀싱", "script-plc.md"],
        ["C: CPLC(실시간) vs capp(백그라운드), C API, CfromScript", "c-programming.md"],
        ["C API 실제 시그니처 (gplib.h / RtGpShm.h pshm)", "c-api.md"],
        ["전체 요소 목록 (firmware intellisense 표)", "firmware/ELEMENTS_INDEX.md"],
        ["함정: 태스크 모델, save/reset, 단위, 모션 안전, 에러 ID", "gotchas.md"],
        ["프로젝트 구조: 폴더, .ppproj, pp_proj.ini 로드 순서", "project-structure.md"],
    ]
    table(slide, rows, y=In(1.55), col_w=[8.2, 3.0], size=12.5, row_h=In(0.5))


def s_mental():
    slide = content("멘탈 모델 — Power PMAC을 다르게 만드는 5가지", "PART 2 · Skill")
    bullets(slide, [
        "**모든 게 명명된 데이터구조 요소** — `Structure[index].Element` (예 `Motor[1].JogSpeed`). 레지스터·I변수 아님. 인덱스는 0부터, **상수/단일 지역변수만**(연산 불가).",
        "**3종 프로그램, 3가지 역할** — 모션(`prog`: 좌표계 경로) / Script PLC(`plc`: 로직·I/O, 1스캔 후 양보) / C(CPLC 실시간 · capp 백그라운드).",
        "**4계층 우선순위** — Phase→Servo→RTI→Background (높을수록 하위 선점). 모션계획·안전체크는 **RTI**.",
        "**좌표계(&) vs 모터(#)** — `&x`는 그룹 경로(run/abort/hold), `#x`는 단일 모터(jog/home/kill). 축정의가 둘을 연결, **축단위 ≠ 모터단위**.",
        "**RAM vs flash** — 다운로드는 RAM만(휘발). `save`=flash 영구화, `$$$`=flash에서 복원, `$$$***`=공장초기화.",
    ], size=16.5, gap=12)


def s_gotchas():
    slide = content("안전 함정(Top gotchas) — 왜 Skill 맨 앞에 박았나", "PART 2 · Skill")
    bullets(slide, [
        "**kill ≠ abort ≠ disable** — kill = 개루프·앰프 off·**감속 없음**(즉시). abort = 폐루프 감속 정지. 수직/중력축은 **지연형**(dkill/ddisable)으로 브레이크 먼저.",
        "**안전장치 대부분 기본 OFF** — 소프트리밋·엔코더손실·앰프폴트는 설정 전엔 꺼짐. **FatalFeLimit** 이 주 폭주 가드 — 0으로 만들지 말 것.",
        "**모터단위 vs 축단위** — jog/home/limit/FatalFeLimit은 **모터단위**. 재스케일하면 모든 한계가 함께 바뀐다.",
        "**실시간 블로킹 금지** — Phase/Servo/RTI/CPLC 코드는 짧고 논블로킹. 초과 시 에러카운터 증가·워치독 트립.",
        "**인덱스는 상수/단일 지역변수** — `Motor[L0+1]`은 불법. 런타임 범위초과는 **에러 없이 손상**.",
    ], size=16, gap=11)
    textbox(slide, MARGIN, In(6.55), CW, In(0.5),
            "→ 그래서 Skill은 **안전 요약을 SKILL.md 최상단**에 두어, 모션·안전 코드 작성 전 항상 점검하게 한다.",
            size=14, color=PRIMARY, bold=True)


def s_skill_use():
    slide = content("Skill 사용 예시 — 그냥 물어보면 된다", "PART 2 · Skill")
    rows = [
        ["이렇게 물으면", "Skill이 하는 일"],
        ["\"1축을 100mm 상대이동 후 복귀하는 모션 프로그램 짜줘\"", "문법에 맞는 prog 생성 (move 모드·축정의)"],
        ["\"이 PLC가 왜 매 스캔마다 안 도는지 봐줘\" (코드 첨부)", "무한 내부루프·스캔모델 위반 진단"],
        ["\"Motor[L0+1].ActPos 는 왜 안 돼?\"", "인덱스에 연산 불가 함정 설명"],
        ["\"capp 와 CPLC 차이, 언제 뭘 써?\"", "실시간 vs 백그라운드 선택 가이드"],
        ["\"근거 reference 가 어디야?\"", "정제·원문 출처(페이지) 제시"],
    ]
    table(slide, rows, y=In(1.7), col_w=[6.6, 4.6], size=13.5, row_h=In(0.62))
    textbox(slide, MARGIN, In(6.1), CW, In(0.5),
            "답변이 실제 요소명·firmware 동작과 일치하면 OK. 미심쩍으면 출처를 되물어 검증한다.",
            size=13, color=GRAY_D)


# ---- Part 3 ----
def s_mcp_what():
    slide = content("MCP 란? — AI가 외부 도구를 호출하는 표준", "PART 3 · MCP")
    bullets(slide, [
        "**Model Context Protocol** — AI(Claude)가 외부 도구·데이터에 **표준 방식**으로 접근하는 개방 프로토콜.",
        "**Claude Code = 클라이언트**, 우리가 만든 **powerpmac-mcp.exe = 서버**.",
        "서버가 \"도구 목록\"을 광고 → Claude가 사용자 의도에 맞는 도구를 골라 호출하고 결과를 받음.",
        "전송 계층: **stdio**(표준입출력)로 **JSON-RPC 2.0** 메시지를 한 줄씩 주고받음.",
        "사용자는 도구 이름을 몰라도 된다 — **자연어 → Claude가 도구 선택 → 실행 전 승인**.",
    ], size=18, gap=13)
    textbox(slide, MARGIN, In(6.4), CW, In(0.5),
            "이 셋(Program · PmacBridge · PdkRuntime)이 어떻게 맞물리는지가 이 파트의 핵심이다.",
            size=13, color=GRAY_D)


def s_arch():
    slide = content("MCP 서버 — 3계층 아키텍처", "PART 3 · MCP")
    box(slide, In(3.0), In(1.55), In(7.3), In(0.7), "Claude Code  (MCP 클라이언트)", fill=PRIMARY, size=15)
    arrow(slide, In(6.65), In(2.25), In(6.65), In(2.75), color=PRIMARY2, width=2.5)
    alabel(slide, In(6.8), In(2.32), In(3.0), "JSON-RPC / stdio", size=11, align=PP_ALIGN.LEFT)
    # container
    cont = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, In(2.2), In(2.8), In(8.9), In(2.55))
    cont.fill.solid(); cont.fill.fore_color.rgb = PANEL
    cont.line.color.rgb = PRIMARY; cont.line.width = Pt(1.5); cont.shadow.inherit = False
    textbox(slide, In(2.35), In(2.88), In(8.6), In(0.35),
            "MCP 서버  powerpmac-mcp.exe  ·  .NET Framework 4.8 · x86", size=12.5, color=PRIMARY, bold=True)
    box(slide, In(2.45), In(3.3), In(8.4), In(0.55),
        ["①  Program.cs — 프로토콜 :  stdin 루프 · JSON-RPC 디스패치 · stdout 보호"],
        fill=MCP_C, size=13, align=PP_ALIGN.LEFT)
    box(slide, In(2.45), In(3.95), In(8.4), In(0.55),
        ["②  PmacBridge.cs — 브리지 :  빌드 · 다운로드 · gpascii/셸 세션 (PDK 래퍼)"],
        fill=TEAL_C, size=13, align=PP_ALIGN.LEFT)
    box(slide, In(2.45), In(4.6), In(8.4), In(0.55),
        ["③  PdkRuntime.cs — 런타임 :  PDK·컴파일러 자동감지 · DLL 동적로딩"],
        fill=RT_C, size=13, align=PP_ALIGN.LEFT)
    arrow(slide, In(4.3), In(5.35), In(4.3), In(5.85), color=RT_C, width=2.2)
    arrow(slide, In(8.7), In(5.35), In(8.7), In(5.85), color=MCP_C, width=2.2)
    box(slide, In(2.3), In(5.85), In(4.0), In(0.8),
        ["PDK DLLs", "ODT.* · cygwin1 · 크로스컴파일러"], fill=RT_C, size=12)
    box(slide, In(6.9), In(5.85), In(4.0), In(0.8),
        ["Power PMAC 컨트롤러", "SSH · gpascii · projpp"], fill=CTRL_C, size=12)
    alabel(slide, In(2.4), In(5.55), In(1.9), "인프로세스", size=10, align=PP_ALIGN.LEFT)
    alabel(slide, In(8.8), In(5.55), In(1.9), "네트워크", size=10, align=PP_ALIGN.LEFT)


def s_proto():
    slide = content("계층 1 — 프로토콜 (Program.cs)", "PART 3 · MCP")
    bullets(slide, [
        "`Main()`이 **PdkRuntime.Init() 먼저** 호출 — 어떤 ODT 타입을 건드리기 전에 PDK 로드 준비.",
        "**stdout은 JSON-RPC 전용** — 진짜 stdout을 잡아두고 `Console.Out`을 stderr로 돌려, 라이브러리의 우발적 print가 프로토콜을 오염시키지 못하게.",
        "stdin을 **한 줄 = 한 메시지**로 읽어 JObject 파싱 → `method` 분기.",
        ("", "initialize / tools/list / tools/call / ping / notifications"),
        "응답 `{jsonrpc, id, result}`, 오류 `{error:{code,message}}`. 한 줄 JSON으로 직렬화해 출력.",
    ], y=In(1.55), h=In(2.7), size=16.5, gap=10)
    code_box(slide, [
        "PdkRuntime.Init();                       // ODT 로드 준비 (최우선)",
        "_out = new StreamWriter(stdout, UTF8(false)){ AutoFlush = true };",
        "Console.SetOut(Console.Error);           // stdout 보호 → stderr",
        "while ((line = stdin.ReadLine()) != null)",
        "    HandleLine(line);                    // JSON-RPC 디스패치",
    ], y=In(4.45), h=In(2.0), size=12.5)


def s_sequence():
    slide = content("MCP 메시지 흐름 (시퀀스)", "PART 3 · MCP")
    cx, sx = In(3.0), In(9.8)
    box(slide, In(1.9), In(1.55), In(2.2), In(0.6), "Claude Code (Client)", fill=PRIMARY, size=12)
    box(slide, In(8.7), In(1.55), In(2.2), In(0.6), "powerpmac-mcp (Server)", fill=MCP_C, size=12)
    # lifelines
    for x in (cx, sx):
        ll = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x, In(2.15), x, In(6.7))
        ll.line.color.rgb = GRAY_D; ll.line.width = Pt(1.0)
        d = ll.line._get_or_add_ln().makeelement(qn('a:prstDash'), {'val': 'dash'})
        ll.line._get_or_add_ln().append(d)
    seq = [
        (2.55, True,  "initialize"),
        (3.05, False, "serverInfo · capabilities {tools}"),
        (3.55, True,  "notifications/initialized"),
        (4.15, True,  "tools/list"),
        (4.65, False, "[ 9개 도구 스키마 ]"),
        (5.30, True,  "tools/call  { name, arguments }"),
        (5.80, False, "content[text]  ·  isError"),
    ]
    for y, l2r, txt in seq:
        if l2r:
            arrow(slide, cx, In(y), sx, In(y), color=PRIMARY2, width=1.75)
        else:
            arrow(slide, sx, In(y), cx, In(y), color=ACCENT, width=1.75)
        alabel(slide, In(3.1), In(y - 0.32), In(6.6), txt, size=11,
               color=(PRIMARY2 if l2r else GRAY_D))
    textbox(slide, MARGIN, In(6.85), CW, In(0.3),
            "→ 파랑=클라이언트 요청, 주황=서버 응답. 모든 교환은 stdio 위의 한 줄 JSON.", size=11, color=GRAY_D)


def s_tools():
    slide = content("도구 9개 — 한눈에", "PART 3 · MCP")
    rows = [
        ["그룹", "도구", "하는 일", "주요 입력(기본값)"],
        ["빌드/다운로드", "build_project", "PC에서 .ppproj 로컬 빌드(컴파일)", "projectPath, configuration(Release)"],
        ["", "download_project", "빌드 결과를 rsync 전송 + projpp 로드", "projectPath, ipAddress, password(deltatau)"],
        ["연결", "connect", "지속 gpascii + 셸 세션 열기", "ipAddress, username(root), port(22)"],
        ["", "disconnect / connection_status", "세션 닫기 / 상태 확인", "—"],
        ["명령", "send_command", "명령 실행(대입·동작)", "command"],
        ["", "get_response / get_responses", "질의 후 응답 / 여러 질의 배치", "command / commands[]"],
        ["셸", "exec_shell", "컨트롤러에서 리눅스 셸 명령", "command (연결 필요)"],
    ]
    gt = table(slide, rows, y=In(1.6), col_w=[1.7, 3.0, 4.0, 4.0], size=12, row_h=In(0.56))
    gt.cell(1, 0).merge(gt.cell(2, 0))   # 빌드/다운로드 그룹 셀 병합
    gt.cell(3, 0).merge(gt.cell(4, 0))   # 연결
    gt.cell(5, 0).merge(gt.cell(6, 0))   # 명령
    textbox(slide, MARGIN, In(6.55), CW, In(0.4),
            "기본 접속값: SSH 포트 22 · 사용자 root · 비밀번호 deltatau (벤더 공장 기본값).",
            size=12, color=GRAY_D)


def s_bridge():
    slide = content("계층 2 — 브리지 (PmacBridge.cs) 개요", "PART 3 · MCP")
    bullets(slide, [
        "**PDK를 얇게 감싼 래퍼** — 헤드리스에서 안전한 작업만 노출.",
        "**lock으로 모든 호출 직렬화** — 단일 세션의 동시성/상태를 보호.",
        "**두 세션을 보유** — gpascii(모션·요소 명령) + terminal(리눅스 셸). connect가 둘 다 엶.",
        "각 도구는 **ToolResult{ Text, IsError }** 를 반환 → Program이 MCP content로 포장.",
        "구성: **① 빌드(로컬)  ② 다운로드(rsync+projpp)  ③ 라이브 세션** — 세 묶음.",
    ], size=18, gap=13)
    box(slide, MARGIN, In(5.6), In(3.85), In(0.95), ["① build_project", "로컬 컴파일"], fill=TEAL_C, size=14)
    box(slide, In(4.74), In(5.6), In(3.85), In(0.95), ["② download_project", "전송 + 로드"], fill=TEAL_C, size=14)
    box(slide, In(8.88), In(5.6), In(3.85), In(0.95), ["③ connect / send / shell", "라이브 세션"], fill=TEAL_C, size=14)


def s_build():
    slide = content("빌드 동작 원리 (build_project)", "PART 3 · MCP")
    bullets(slide, [
        "PDK의 **Build.BuildProject(.ppproj, config)** 를 **인프로세스**로 호출.",
        "이 한 번이 **Script 컴파일 + ARM C 크로스컴파일**을 모두 수행 → `Bin/<config>/*.out` 생성.",
        "**PPMAC460CompileTask.dll**(+의존)을 exe 옆에 복사 → MSBuild 태스크가 BaseDirectory에서 로드.",
        "결과의 **TotalErrors / TotalWarnings** 와 에러·경고 목록을 파싱해 텍스트로 보고.",
        "C 코드를 바꿨으면 **빌드 먼저**, 그다음 다운로드 (다운로드는 C를 재빌드하지 않음).",
    ], y=In(1.55), h=In(3.0), size=17, gap=11)
    code_box(slide, [
        "var build = new Build();",
        "Build.BuildResults r = build.BuildProject(projectPath, config);",
        "bool ok = r.TotalErrors == 0;     //  Script + ARM C 한 번에",
        "//  -> \"Build SUCCEEDED/FAILED (Release)  Errors: 0  Warnings: 0\"",
    ], y=In(4.7), h=In(1.7), size=12.5)


def s_download():
    slide = content("다운로드 동작 원리 — 헤드리스 PTY 우회", "PART 3 · MCP")
    textbox(slide, MARGIN, In(1.5), CW, In(0.4),
            "문제: PDK의 RSYNC 다운로드는 cygwin ssh가 대화형 **PTY**를 요구 → 헤드리스 MCP엔 PTY가 없어 **무한 대기(행)**.",
            size=13.5, color=RED_C, bold=True)
    # 막힌 경로
    box(slide, In(0.7), In(2.1), In(2.6), In(0.85), ["MCP (헤드리스)", "stdio 프로세스"], fill=CTRL_C, size=12)
    arrow(slide, In(3.3), In(2.52), In(4.3), In(2.52), color=RED_C, width=2.0, dashed=True)
    box(slide, In(4.3), In(2.1), In(3.5), In(0.85), ["PDK RSYNC 호출", "PTY 없음 → 행(hang)"],
        fill=RED_C, size=12)
    alabel(slide, In(3.2), In(2.18), In(1.2), "사용 안 함", size=10, color=RED_C)
    # 해결 경로
    textbox(slide, MARGIN, In(3.2), CW, In(0.35), "해결 — 직접 배치 + 새 콘솔로 PTY 확보:", size=14, color=SKILL_C, bold=True)
    box(slide, In(0.7), In(3.7), In(2.7), In(1.15),
        ["① 배치(.cmd) 생성", "sshpass+rsync → projpp"], fill=MCP_C, size=12.5)
    arrow(slide, In(3.4), In(4.27), In(3.85), In(4.27), color=GRAY_D)
    box(slide, In(3.85), In(3.6), In(3.0), In(1.35),
        ["② cmd.exe 실행", "**UseShellExecute=true**", "→ 새 콘솔 할당", "→ cygwin ssh에 PTY"], fill=SKILL_C, size=12)
    arrow(slide, In(6.9), In(4.27), In(7.35), In(4.27), color=GRAY_D)
    box(slide, In(7.35), In(3.7), In(2.7), In(1.15),
        ["③ rsync 전송", "**상대경로 ./**", "(`C:\\` 절대경로는 원격 오해)"], fill=TEAL_C, size=12)
    arrow(slide, In(10.1), In(4.27), In(10.55), In(4.27), color=GRAY_D)
    box(slide, In(10.55), In(3.7), In(2.15), In(1.15),
        ["④ 컨트롤러", "projpp -l", "PMAC 버퍼 로드"], fill=CTRL_C, size=12)
    bullets(slide, [
        "성패 판정: 출력의 **RSYNC_EXIT=0** 과 **PROJPP_EXIT=0 / projpp errors=0** 확인.",
        "실증: 전체 프로젝트를 헤드리스로 **약 8초**에 전송·로드.",
    ], y=In(5.2), size=14, gap=8)
    box(slide, MARGIN, In(6.12), In(12.13), In(0.66),
        ["**PTY** (pseudo-terminal, 의사 터미널): ssh 같은 대화형 프로그램이 비밀번호 입력·터미널 "
         "제어를 위해 요구하는 가상 터미널. 헤드리스(백그라운드) 프로세스엔 없어, 그대로 호출하면 ssh가 멈춘다."],
        fill=PANEL, txt=INK, size=11.5, bold=False, align=PP_ALIGN.LEFT, line=PANEL2)


def s_live():
    slide = content("라이브 세션 (connect · get · send · shell)", "PART 3 · MCP")
    bullets(slide, [
        "**connect(ip, user=root, pw=deltatau, port=22)** — gpascii 세션 + (best-effort) 셸 터미널을 함께 엶.",
        "connect 동안만 **NativeSearch(true)** — PDK 네이티브 DLL(DKeyLib/EncPass/라이선스) 검색을 일시 활성(끝나면 off).",
        "**get_response** (질의 후 응답) · **send_command** (대입·동작, 응답에 ERR 있으면 오류로) · **get_responses** (여러 질의를 **1왕복 배치**).",
        "**exec_shell** — terminal로 리눅스 명령(ls, 로그 확인 등).",
        "**세션은 한 번에 하나** — 새 connect가 이전 세션을 대체. 모든 호출은 lock으로 직렬화.",
    ], size=17, gap=12)
    textbox(slide, MARGIN, In(6.4), CW, In(0.5),
            "예: \"연결하고 Motor[1]~Motor[4] ActPos 한 번에 읽어줘\" → connect + get_responses (1왕복).",
            size=13, color=GRAY_D)


def s_runtime():
    slide = content("계층 3 — PDK 런타임 (PdkRuntime.cs)", "PART 3 · MCP")
    bullets(slide, [
        "**자동 감지로 이식성** — 환경변수 없이도 머신마다 PDK·컴파일러를 찾는다.",
        (1, "PDK: `POWERPMAC_PDK_HOME` → 레지스트리 \"PowerPMAC Development Kit\" → `…\\PowerPMAC\\<버전>\\PDK` → 공통경로"),
        (1, "컴파일러: `POWERPMAC_COMPILERS_HOME` → `DTBUILDPATH` → 레지스트리 → `C:\\DeltaTau\\PowerPMAC\\Compilers`"),
        "**AssemblyResolve** — `ODT.*` 등 매니지드 DLL을 PDK 폴더에서 **동적 로드**(출력 폴더에 복사하지 않아도 됨).",
        "**NativeSearch는 connect 동안만** — `SetDllDirectory(PDK)`를 상시 켜면 C 빌드가 엉뚱한 cygwin1.dll을 물어 깨진다. 그래서 **일시 적용**.",
    ], size=16.5, gap=11)
    textbox(slide, MARGIN, In(6.45), CW, In(0.5),
            "→ 핵심: \"제자리(in-place)\"의 PDK 런타임을 복사 없이 로드 가능하게 만들고, 머신별 경로를 스스로 찾는다.",
            size=13, color=PRIMARY, bold=True)


def s_pdkdetect():
    slide = content("PDK 판정 로직 — 무엇으로 PDK를 식별하나", "PART 3 · MCP")
    bullets(slide, [
        "PDK 폴더는 **DLL 3종**으로 식별:",
        (1, "ODT.PowerPmacBuildAndDownload.dll  ·  PPMAC460CompileTask.dll  ·  cygwin1.dll"),
        "**CLLLicFile.lic 는 요구하지 않는다** — 라이선스의 정식 위치는 `System32`/`SysWOW64`라 PDK 폴더엔 없을 수 있음(체험판 등). 요구하면 정상 PDK를 **false-reject**.",
        "라이선스 유효성은 **빌드 시 컴파일러가 강제** — 탐지 단계에서 설치를 막지 않는다.",
        "최근 수정(커밋 **9cc5f8d**)으로 동료 PC의 \"PDK 못 찾음\" 오탐을 해결.",
    ], y=In(1.55), h=In(3.0), size=16.5, gap=11)
    code_box(slide, [
        "private static bool IsPdk(string p) =>",
        "    Directory.Exists(p)",
        "    && File.Exists(Path.Combine(p, \"ODT.PowerPmacBuildAndDownload.dll\"))",
        "    && File.Exists(Path.Combine(p, \"PPMAC460CompileTask.dll\"))",
        "    && File.Exists(Path.Combine(p, \"cygwin1.dll\"));   //  .lic 비요구",
    ], y=In(4.7), h=In(1.9), size=12)


def s_x86():
    slide = content("왜 net48 + x86 인가", "PART 3 · MCP")
    bullets(slide, [
        "PDK 런타임이 **32비트**다 — cygwin1.dll · ODT.* · PPMAC460CompileTask · Renci.SshNet.",
        "매니지드 DLL은 **.NET Framework + WinForms** 대상.",
        "→ 호스트 프로세스도 **net48 · x86** 이어야 네이티브 로드가 성공한다.",
        (1, "csproj: PlatformTarget=x86, Prefer32Bit=true, TargetFramework=net48"),
        "불일치하면 **BadImageFormatException** — 64비트 호스트는 32비트 네이티브 DLL을 못 올린다.",
    ], size=17.5, gap=13)
    box(slide, MARGIN, In(5.7), In(12.1), In(0.95),
        ["호스트(net48·x86)  ⇄  매니지드 ODT.*  ⇄  네이티브 cygwin1/DKeyLib  —  **세 비트수가 일치해야 로드**"],
        fill=PRIMARY, size=14, rounded=True)


# ---- Part 4 ----
def s_setup():
    slide = content("설치 자동화 — setup.ps1", "PART 4 · 설치·배포")
    bullets(slide, [
        "**① 콘솔 UTF-8** — dotnet의 한글 출력이 CP949 콘솔에서 깨지지 않게 `[Console]::OutputEncoding = UTF8`.",
        "**② PDK 감지** — `Test-IsPdk`(DLL 3종). 못 찾으면 `-PdkHome \"C:\\...\\PDK\"` 로 지정.",
        "**③ dotnet 빌드** — `/p:PdkHome=<감지경로>` 로 powerpmac-mcp.exe(x86) 생성.",
        "**④ Skill 설치** — `~/.claude/skills/powerpmac-dev` 연결.",
        "**⑤ MCP 등록** — `claude` CLI가 PATH에 있으면 `claude mcp add`, 없으면(데스크톱 앱) **~/.claude.json 직접 편집**(Register-McpDirect: 백업·멱등·stdio 스키마 삽입).",
        "**-SkillOnly** — PDK 없는 PC는 Skill만 설치(코드 작성·리뷰 전용).",
    ], size=16, gap=10)
    textbox(slide, MARGIN, In(6.55), CW, In(0.4),
            "한 번 실행 → 감지·빌드·등록 자동. 끝나면 Claude Code 재시작.", size=13, color=GRAY_D)


def s_deploy():
    slide = content("배포 모델 — 듀얼 리모트", "PART 4 · 설치·배포")
    box(slide, In(5.3), In(1.7), In(2.7), In(1.0), ["로컬 작업본", "git push (1회)"], fill=PRIMARY, size=14)
    arrow(slide, In(6.0), In(2.7), In(3.4), In(3.5), color=GRAY_D, width=2.2)
    arrow(slide, In(7.3), In(2.7), In(9.9), In(3.5), color=GRAY_D, width=2.2)
    box(slide, In(1.3), In(3.5), In(4.3), In(1.1),
        ["사내 Z: 공유 (bare)", "내부 백업·배포"], fill=CTRL_C, size=13)
    box(slide, In(7.7), In(3.5), In(4.3), In(1.1),
        ["GitHub (Public)", "minimin333/PowerPMAC_MCP"], fill=MCP_C, size=13)
    bullets(slide, [
        "origin에 push URL 2개를 등록(set-url --add --push) → **git push 한 번에 양쪽 동시 반영**.",
        "**bin/ 은 gitignore** — exe는 배포하지 않고, 동료가 setup.ps1로 **각자 빌드**(PDK 경로가 PC마다 다름).",
        "공개 전 **사내 경로(Z:) 스크럽** — 공개 문서엔 GitHub clone URL만 남김.",
        "업데이트: 동료는 **git pull + setup.ps1 재실행**.",
    ], y=In(4.9), size=15, gap=9)


def s_workflow():
    slide = content("전체 개발 워크플로", "PART 4 · 설치·배포")
    # top row L->R
    box(slide, In(0.6), In(1.7), In(2.0), In(0.95), ["① 코드 작성·수정", "(Skill)"], fill=SKILL_C, size=12.5)
    arrow(slide, In(2.6), In(2.17), In(3.0), In(2.17))
    box(slide, In(3.0), In(1.78), In(1.85), In(0.8), "build_project", fill=MCP_C, size=12.5)
    arrow(slide, In(4.85), In(2.17), In(5.25), In(2.17))
    diamond(slide, In(5.25), In(1.55), In(1.7), In(1.25), ["빌드", "성공?"], fill=ACCENT, size=12)
    arrow(slide, In(6.95), In(2.17), In(7.45), In(2.17))
    alabel(slide, In(6.85), In(1.82), In(0.7), "예", size=10, color=SKILL_C)
    box(slide, In(7.45), In(1.78), In(2.3), In(0.8), ["download_project", "(→ RAM)"], fill=MCP_C, size=12)
    arrow(slide, In(9.75), In(2.17), In(10.2), In(2.17))
    box(slide, In(10.2), In(1.78), In(2.0), In(0.8), ["connect", "(gpascii)"], fill=MCP_C, size=12)
    # down on right
    arrow(slide, In(11.2), In(2.58), In(11.2), In(3.55), width=2.2)
    # bottom row R->L
    box(slide, In(10.0), In(3.55), In(2.3), In(1.0), ["get_response /", "send_command (검증)"], fill=TEAL_C, size=12)
    arrow(slide, In(10.0), In(4.05), In(9.55), In(4.05))
    diamond(slide, In(7.85), In(3.45), In(1.7), In(1.2), ["정상", "동작?"], fill=ACCENT, size=12)
    arrow(slide, In(7.85), In(4.05), In(7.4), In(4.05))
    alabel(slide, In(7.3), In(3.7), In(0.7), "예", size=10, color=SKILL_C)
    box(slide, In(5.2), In(3.6), In(2.2), In(0.95), ["send_command", "\"save\" (flash)"], fill=MCP_C, size=12)
    arrow(slide, In(5.2), In(4.07), In(4.75), In(4.07))
    box(slide, In(2.9), In(3.65), In(1.85), In(0.85), "완료", fill=PRIMARY, size=14)
    # feedback loops — ㄷ자 경로(아래로 우회)로 되돌아가는 방향을 명확히
    line(slide, In(6.1), In(2.8), In(6.1), In(3.08), color=RED_C, width=1.5, dashed=True)
    line(slide, In(6.1), In(3.08), In(1.6), In(3.08), color=RED_C, width=1.5, dashed=True)
    arrow(slide, In(1.6), In(3.08), In(1.6), In(2.66), color=RED_C, width=1.5, dashed=True)
    alabel(slide, In(2.55), In(2.82), In(3.2), "빌드 실패 → 수정", size=10, color=RED_C, align=PP_ALIGN.LEFT)
    line(slide, In(8.7), In(4.65), In(8.7), In(5.22), color=RED_C, width=1.5, dashed=True)
    line(slide, In(8.7), In(5.22), In(1.3), In(5.22), color=RED_C, width=1.5, dashed=True)
    arrow(slide, In(1.3), In(5.22), In(1.3), In(2.66), color=RED_C, width=1.5, dashed=True)
    alabel(slide, In(2.7), In(4.98), In(3.6), "검증 실패 → ①로 수정", size=10, color=RED_C, align=PP_ALIGN.LEFT)
    textbox(slide, MARGIN, In(6.2), CW, In(0.7),
            ["다운로드는 RAM(휘발) — 영구 저장은 컨트롤러 save 필요. \"save 해줘\" → send_command \"save\".",
             "사용자는 자연어로 말하고, Claude가 단계마다 알맞은 도구를 골라 승인을 받고 실행한다."],
            size=12.5, color=GRAY_D, gap=3)


def s_proof():
    slide = content("실증 사례 — 실물 컨트롤러 검증", "PART 4 · 설치·배포")
    bullets(slide, [
        "대상: 실물 **192.168.0.200**, 펌웨어 **2.8.3.0**.",
        "**build_project** — ARM C 크로스컴파일 포함 빌드 성공.",
        "**download_project** — 전체 프로젝트를 헤드리스로 **~8초**에 전송·로드. 이후 `list plc 0` / `list prog 2`가 방금 코드를 그대로 반환.",
        "**connect + get_responses** — 라이브 요소 일괄 조회(1왕복). **exec_shell** 로 셸 확인.",
        "실제 **Motor 1 조그 → 킬**까지 안전 점검(AmpEna·FeFatal·InPos)과 함께 데모 완료.",
    ], size=17, gap=12)
    box(slide, MARGIN, In(5.9), In(12.1), In(0.9),
        ["작성(Skill) → build → download(~8초) → connect → get/send → save  —  한 대화 안에서 전체 루프 검증됨"],
        fill=SKILL_C, size=13.5)


# ---- Part 5 ----
def s_decisions():
    slide = content("핵심 설계 결정 요약", "PART 5 · 마무리")
    rows = [
        ["결정", "왜 / 효과"],
        ["stdout 프로토콜 전용 (Console.Out→stderr)", "라이브러리 print가 JSON-RPC를 오염시키지 못하게"],
        ["헤드리스 PTY 우회 (UseShellExecute 새 콘솔)", "cygwin ssh에 PTY 제공 — 다운로드 행(hang) 방지"],
        ["rsync 상대경로 (./ )", "`C:\\` 의 드라이브 : 를 원격호스트로 오해하는 문제 회피"],
        ["NativeSearch 일시 적용 (connect 동안만)", "SetDllDirectory 상시 적용이 C 빌드를 깨는 것 방지"],
        ["PDK·컴파일러 자동 감지", "환경변수 없이도 머신마다 동작 — 무설정 이식성"],
        ["net48 · x86 고정", "32비트 PDK 네이티브 DLL 로드 보장"],
        ["PDK 탐지에서 .lic 비요구", "체험판·라이선스 위치차로 인한 false-reject 방지"],
    ]
    table(slide, rows, y=In(1.55), col_w=[5.4, 6.0], size=13, row_h=In(0.56))
    textbox(slide, MARGIN, In(6.7), CW, In(0.4),
            "공통 주제: 헤드리스 · 이식성 · 안전.", size=14, color=PRIMARY, bold=True)


def s_closing():
    slide = content("마무리 · 참고자료", "PART 5 · 마무리")
    bullets(slide, [
        "**사용자 매뉴얼** — docs/manual.md · manual.html · manual.pdf (설치→등록→사용 한 번에)",
        "**저장소 개요** — README.md   /   **설치·문제해결** — INSTALL.md",
        "**MCP 내부 동작** — mcp-server/README.md (헤드리스 PTY · rsync + projpp)",
        "**Skill 지식 구조** — Skills/powerpmac-dev/SKILL.md (+ reference/)",
        "**공개 저장소** — github.com/minimin333/PowerPMAC_MCP",
    ], y=In(1.7), size=17, gap=12)
    box(slide, MARGIN, In(5.5), In(12.1), In(1.2),
        ["기억할 한 가지", "\"물어보면 된다 — 자연어로 작성 · 빌드 · 다운로드 · 제어까지.\""],
        fill=PRIMARY, size=18)


# ================================================================ 발표 노트(스크립트)
def set_notes(slide, text):
    """슬라이드 노트(발표 스크립트) 주입 — 슬라이드쇼엔 안 보이고 발표자 노트에만 표시."""
    tf = slide.notes_slide.notes_text_frame
    tf.text = text
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.size = Pt(13)
            set_font(run, FONT_KR)


# 슬라이드 순서(= build()의 호출 순서)에 1:1로 대응하는 발표 스크립트
NOTES = [
    # 1 표지
    """OMRON Delta Tau Power PMAC 모션 컨트롤러 개발을 돕는 도구의 구조와 동작 원리를 소개합니다. 핵심은 두 가지입니다 — Claude Code에 더해지는 powerpmac-dev Skill(지식)과 powerpmac MCP 서버(조작). 이 둘이 어떻게 자연어만으로 코드 작성부터 실제 장비 제어까지 이어주는지 차례로 살펴보겠습니다.""",
    # 2 목차
    """전체는 5개 파트입니다. Part 1은 이 도구가 푸는 문제와 전체 그림, Part 2는 Skill이 Power PMAC 지식을 제공하는 방식, Part 3은 MCP 서버가 실제 컨트롤러를 조작하는 동작 원리로 이 발표의 핵심입니다. Part 4는 설치와 배포, Part 5는 핵심 설계 결정을 정리합니다. 시간이 부족하면 Part 3에 집중하시면 됩니다.""",
    # 3 divider Part1
    """먼저 이 도구가 어떤 불편을 해결하고, 전체 구성이 어떻게 생겼는지부터 봅니다.""",
    # 4 문제의식
    """기존 Power PMAC 개발의 마찰부터 짚습니다. 첫째, 매뉴얼이 방대합니다 — 공식 문서가 4종 약 3,260쪽이고 요소·명령 이름이 펌웨어 버전마다 달라 매번 찾아야 합니다. 둘째, IDE에서 빌드·다운로드·터미널 조회를 사람이 손으로 반복합니다. 셋째, kill과 abort의 차이, 기본으로 꺼진 안전장치, 모터·축 단위 혼동 같은 함정이 자칫 장비 사고로 이어집니다. 그래서 지식과 조작을 모두 거들어 주는 AI 에이전트가 있으면 좋겠다는 게 출발점입니다.""",
    # 5 해결책
    """해결책은 Claude Code라는 AI 코딩 에이전트에 두 가지를 더하는 것입니다. Skill은 Power PMAC을 '이해'하게 해서 정확한 코드 작성과 리뷰를 돕고, MCP는 실제 컨트롤러를 '조작'해 빌드·다운로드·조회를 자연어로 처리합니다. 오른쪽 Before/After를 보면, 예전엔 IDE를 열고 수동으로 여러 도구를 오갔다면 이제는 '이 프로젝트 빌드해서 200번 장비에 다운로드하고 Motor 1 위치 읽어줘' 한 마디로 전체 흐름이 진행됩니다.""",
    # 6 두 축
    """Skill과 MCP의 역할을 명확히 구분합니다. Skill은 지식입니다 — 마크다운 지식베이스로, 질문하면 자동 적용되고 PDK 없이도 씁니다. MCP는 조작입니다 — 실행 파일 서버로, '해줘'라고 하면 도구를 호출하며 빌드·다운로드에 PDK가 필요합니다. 핵심은 둘이 상보적이라는 점입니다. Skill로 코드를 작성하고, MCP로 그 코드를 장비에 돌립니다.""",
    # 7 시스템 구성도
    """전체 그림입니다. 개발자가 자연어로 지시하면 Claude Code가 받습니다. 지식이 필요하면 Skill을, 장비 조작이 필요하면 MCP 서버를 부릅니다. 가운데 초록 화살표가 중요한데, Skill로 작성·리뷰한 코드를 MCP 서버가 빌드·다운로드해 실제 컨트롤러에 반영하는 흐름입니다. MCP 서버는 PDK의 DLL과 크로스컴파일러를 인프로세스로 쓰고, SSH·gpascii로 컨트롤러와 통신합니다. 사용자는 어떤 도구가 호출되는지 몰라도 됩니다.""",
    # 8 divider Part2
    """이제 Skill이 어떻게 Power PMAC을 '이해'하게 만드는지 봅니다.""",
    # 9 Skill 메커니즘
    """Skill은 Claude Code의 지식 모듈입니다. 질문이 Power PMAC 도메인에 맞으면 자동으로 로드되며, 사용자가 따로 부를 필요가 없습니다. 설치는 홈 폴더의 skills 디렉터리에 연결되는 형태입니다. 핵심은 점진적 로드입니다 — 항상 읽는 건 지도 역할의 SKILL.md 한 장뿐이고, 깊은 내용은 필요할 때만 reference 폴더의 해당 파일을 펼쳐 토큰을 아낍니다.""",
    # 10 Skill 구조
    """Skill은 세 층으로 정리돼 있습니다. SKILL.md는 지도와 안전 요약 — 멘탈모델, 핵심 함정, 라우팅 표가 들어 있어 항상 읽힙니다. reference 폴더는 도메인별로 정제한 distillation이고, 그 아래 raw 폴더에는 공식 매뉴얼 원문 약 3,260쪽이 검색용으로 들어 있습니다. 동작 규율이 중요한데, 요소 이름이 불확실하면 원문을 직접 검색해 확인하고 절대 지어내지 않으며, 비자명한 규칙은 매뉴얼 페이지를 인용합니다.""",
    # 11 라우팅
    """이 표는 Skill의 지식 라우팅입니다. 질문 도메인에 따라 어떤 reference 파일을 펼칠지 매핑돼 있습니다. 문법·변수는 syntax-rules, 데이터구조는 data-structure, 모션 프로그램은 script-motion, PLC는 script-plc, C는 c-programming과 c-api, 전체 요소 목록은 firmware 인덱스, 함정은 gotchas로 갑니다. 덕분에 방대한 원문을 다 읽지 않고도 필요한 부분만 정확히 참조합니다.""",
    # 12 멘탈모델
    """Power PMAC을 다른 컨트롤러와 다르게 만드는 다섯 가지입니다. 첫째, 모든 게 명명된 데이터구조 요소입니다 — 레지스터가 아니라 Motor 1의 JogSpeed 같은 이름으로 접근하고, 인덱스엔 상수나 단일 지역변수만 씁니다. 둘째, 프로그램이 세 종류로 좌표계 경로용 모션, 로직용 PLC, 실시간·백그라운드 C입니다. 셋째, Phase·Servo·RTI·Background 네 우선순위 계층이 있습니다. 넷째, 좌표계 명령과 모터 명령이 구분됩니다. 다섯째, 다운로드는 RAM에만 올라가고 save를 해야 flash에 영구 저장됩니다. 이 다섯이 머릿속에 있어야 코드가 보입니다.""",
    # 13 gotchas
    """안전 함정을 Skill이 왜 맨 앞에 두는지 보여줍니다. kill은 즉시 개루프로 감속 없이 멈추고 abort는 폐루프로 감속 정지라, 수직축에선 브레이크가 먼저 걸리는 지연형을 써야 합니다. 소프트리밋·엔코더손실·앰프폴트 같은 안전장치는 대부분 기본으로 꺼져 있고, FatalFeLimit이 주된 폭주 방지선입니다. jog·home·리밋은 모터 단위라 재스케일하면 모든 한계가 함께 바뀝니다. 실시간 코드는 블로킹하면 안 되고, 인덱스에는 연산을 못 씁니다. 그래서 Skill은 모션·안전 코드를 내놓기 전에 이 요약을 항상 점검합니다.""",
    # 14 Skill 사용 예시
    """실제로는 그냥 물어보면 됩니다. '1축을 100mm 상대 이동 후 복귀하는 모션 프로그램 짜줘' 하면 문법에 맞는 prog를 만들고, 'PLC가 왜 매 스캔 안 도는지 봐줘' 하면 스캔 모델 위반을 진단합니다. 인덱스에 연산을 쓴 코드는 왜 안 되는지 설명하고, capp와 CPLC 차이도 안내합니다. 답이 미심쩍으면 '근거 reference가 어디야'라고 되물어 출처를 확인할 수 있습니다.""",
    # 15 divider Part3
    """여기가 핵심입니다. MCP 서버가 실제 컨트롤러를 어떻게 조작하는지, 내부 동작 원리를 봅니다.""",
    # 16 MCP란
    """MCP는 Model Context Protocol의 약자로, AI가 외부 도구에 표준 방식으로 접근하는 개방 프로토콜입니다. 여기서 Claude Code가 클라이언트, 우리가 만든 powerpmac-mcp.exe가 서버입니다. 서버가 제공 도구 목록을 알리면 Claude가 사용자 의도에 맞는 도구를 골라 호출합니다. 전송은 표준입출력 위의 JSON-RPC이고, 사용자는 도구 이름을 몰라도 자연어로 말하면 Claude가 선택해 실행 전 승인을 받습니다.""",
    # 17 3계층 아키텍처
    """MCP 서버 내부는 세 계층입니다. 위에서부터 Program.cs가 프로토콜로 표준입출력 JSON-RPC를 주고받고, PmacBridge.cs가 브리지로 빌드·다운로드·세션 같은 실제 작업을 PDK를 감싸 수행하며, PdkRuntime.cs가 런타임으로 PDK와 컴파일러를 자동으로 찾아 DLL을 로딩합니다. 이 서버는 PDK DLL을 인프로세스로 올리고 네트워크로 컨트롤러와 통신합니다. 다음 슬라이드부터 각 계층을 차례로 봅니다.""",
    # 18 계층1 프로토콜
    """첫 계층, 프로토콜입니다. Main 함수가 가장 먼저 PdkRuntime.Init을 호출해 PDK 로딩을 준비합니다. 중요한 트릭이 stdout 보호입니다 — 표준출력은 JSON-RPC 전용이라, 진짜 출력 핸들을 따로 잡아두고 Console.Out은 표준에러로 돌립니다. 라이브러리가 실수로 출력해도 프로토콜이 깨지지 않게 하려는 겁니다. 표준입력은 한 줄이 한 메시지라, 읽어서 파싱하고 method에 따라 분기합니다.""",
    # 19 시퀀스
    """클라이언트와 서버가 주고받는 순서입니다. 먼저 initialize로 서로 능력을 교환하고 initialized 알림을 보냅니다. 그다음 tools/list로 서버가 9개 도구의 스키마를 알려주고, 실제 작업은 tools/call로 도구 이름과 인자를 넘기면 서버가 결과를 content로 돌려줍니다. 파란 화살표가 클라이언트 요청, 주황이 서버 응답이며, 모든 교환은 표준입출력 위의 한 줄 JSON입니다.""",
    # 20 도구 9개
    """서버가 제공하는 도구는 9개입니다. 빌드·다운로드 그룹은 로컬 빌드와 컨트롤러 전송, 연결 그룹은 세션 열기·닫기·상태 확인, 명령 그룹은 명령 실행과 조회·배치 조회, 셸 그룹은 리눅스 명령입니다. 기본 접속값은 SSH 22번 포트, 사용자 root, 비밀번호 deltatau입니다. 사용자는 이 이름을 외울 필요 없이 자연어로 말하면 Claude가 골라 씁니다.""",
    # 21 계층2 브리지
    """둘째 계층, 브리지입니다. PDK를 얇게 감싼 래퍼로, 헤드리스 환경에서 안전한 작업만 노출합니다. lock으로 모든 호출을 직렬화해 단일 세션 상태를 보호하고, gpascii 세션과 리눅스 셸 터미널 두 개를 함께 엽니다. 각 도구는 결과 텍스트와 에러 플래그를 담은 ToolResult를 돌려주고, Program이 이를 MCP content로 포장합니다. 기능은 빌드, 다운로드, 라이브 세션 세 묶음입니다.""",
    # 22 빌드
    """빌드는 PDK의 BuildProject를 인프로세스로 한 번 호출하면, Script 컴파일과 ARM C 크로스컴파일이 모두 수행되어 바이너리가 생성됩니다. 이를 위해 컴파일 태스크 DLL을 실행 파일 옆에 복사해 MSBuild가 찾게 합니다. 결과의 에러·경고 개수와 목록을 파싱해 텍스트로 보고합니다. 주의할 점은 C 코드를 바꿨으면 반드시 빌드를 먼저 하고 다운로드해야 한다는 것입니다 — 다운로드는 C를 다시 빌드하지 않습니다.""",
    # 23 다운로드
    """다운로드는 이 도구에서 가장 까다로웠던 부분입니다. PDK의 기본 RSYNC 다운로드는 cygwin ssh가 대화형 PTY를 요구하는데, 헤드리스 MCP에는 PTY가 없어 그냥 멈춰버립니다. 그래서 세 가지 트릭을 씁니다. 첫째 sshpass·rsync·projpp를 실행하는 배치 파일을 직접 만들고, 둘째 cmd.exe를 UseShellExecute로 띄워 윈도우가 새 콘솔을 할당하게 해 ssh에 PTY를 줍니다. 셋째 rsync 소스를 상대경로로 줘서 C 드라이브 콜론을 원격 호스트로 오해하는 문제를 피합니다. 마지막에 컨트롤러에서 projpp가 로드하고, 전체가 약 8초에 끝납니다. 아래 PTY 설명도 참고하세요.""",
    # 24 라이브 세션
    """라이브 세션입니다. connect는 IP만 주면 기본값으로 gpascii 세션과 셸 터미널을 함께 엽니다. 연결하는 동안만 NativeSearch를 켜서 PDK 네이티브 DLL을 찾고 끝나면 끕니다. get_response는 값을 조회하고, send_command는 대입이나 동작을 실행하며, get_responses는 여러 질의를 한 번의 왕복으로 처리합니다. exec_shell로 리눅스 명령도 됩니다. 세션은 한 번에 하나라, 새로 연결하면 이전 세션을 대체합니다.""",
    # 25 계층3 런타임
    """셋째 계층, 런타임입니다. 이식성의 핵심으로, 환경변수 없이도 머신마다 PDK와 컴파일러를 자동으로 찾습니다. AssemblyResolve로 ODT 같은 매니지드 DLL을 PDK 폴더에서 동적 로딩해, 출력 폴더에 복사하지 않아도 됩니다. 한 가지 미묘한 점은 NativeSearch를 연결하는 동안만 켠다는 것입니다 — SetDllDirectory를 계속 켜두면 C 빌드가 엉뚱한 cygwin DLL을 물어 깨지기 때문에 일시적으로만 적용합니다.""",
    # 26 PDK 판정
    """PDK 폴더를 무엇으로 식별하는지입니다. DLL 세 개 — 빌드앤다운로드, 컴파일태스크, cygwin1 — 이 있으면 PDK로 봅니다. 예전엔 라이선스 파일도 요구했는데, 그 파일의 정식 위치가 시스템 폴더라 PDK 폴더엔 없을 수 있어서 체험판 등에서 멀쩡한 PDK를 잘못 거부하는 문제가 있었습니다. 그래서 라이선스 요구를 빼고 DLL로만 판정하도록 최근에 고쳤습니다. 라이선스 유효성은 어차피 빌드할 때 컴파일러가 확인합니다.""",
    # 27 net48 x86
    """왜 닷넷 프레임워크 4.8에 32비트인가입니다. PDK 런타임의 cygwin1, ODT, 컴파일태스크 등이 모두 32비트이고 매니지드 DLL은 닷넷 프레임워크와 윈폼 대상이라, 호스트 프로세스도 같은 32비트여야 네이티브 로드가 됩니다. 비트 수가 안 맞으면 BadImageFormatException이 나며 64비트 호스트는 32비트 DLL을 못 올립니다. 그래서 csproj에서 x86으로 고정합니다.""",
    # 28 divider Part4
    """이제 설치와 배포, 그리고 전체 워크플로와 실증 사례를 봅니다.""",
    # 29 setup
    """설치는 setup.ps1 한 번으로 끝납니다. 먼저 콘솔을 UTF-8로 맞춰 dotnet의 한글 출력이 깨지지 않게 하고, PDK를 감지한 뒤 그 경로로 dotnet 빌드를 해 실행 파일을 만듭니다. Skill을 홈 폴더에 연결하고 MCP를 등록하는데, claude CLI가 PATH에 있으면 명령으로, 없으면 데스크톱 앱 환경이라 보고 설정 파일을 직접 편집합니다. PDK가 없는 PC는 SkillOnly 옵션으로 코드 작성·리뷰만 쓸 수 있습니다.""",
    # 30 배포
    """배포는 듀얼 리모트입니다. origin에 푸시 URL을 두 개 등록해서, git push 한 번이면 사내 공유와 공개 GitHub 양쪽에 동시에 반영됩니다. 빌드 산출물 bin 폴더는 git에서 제외하는데, PDK 경로가 PC마다 달라 각자 빌드해야 하기 때문입니다. 공개 전에는 사내 경로를 모두 제거했고, 동료는 git pull 후 setup을 다시 돌리면 업데이트됩니다.""",
    # 31 워크플로
    """전체 개발 흐름입니다. Skill로 코드를 작성·수정하고 build_project로 빌드합니다. 빌드가 실패하면 다시 작성으로 돌아갑니다. 성공하면 download_project로 RAM에 올리고, connect로 연결해 값을 조회하거나 명령을 보내 검증합니다. 정상이 아니면 다시 작성으로 돌아가고, 정상이면 save로 flash에 영구 저장하고 끝납니다. 빨간 점선이 실패 시 되돌아가는 경로입니다. 사용자는 자연어로 말하고 Claude가 단계마다 도구를 골라 승인받고 실행합니다.""",
    # 32 실증
    """실제 장비로 검증한 결과입니다. IP 200번, 펌웨어 2.8.3.0에서 빌드는 ARM C 크로스컴파일 포함해 성공했고, 다운로드는 전체 프로젝트를 헤드리스로 약 8초에 전송·로드했습니다. 이후 list 명령으로 방금 올린 코드가 그대로 돌아오는 것을 확인했고, 연결해서 요소를 일괄 조회하고 셸도 실행했습니다. 실제로 모터 1을 조그하고 킬하는 것까지 안전 점검과 함께 시연했습니다. 한 대화 안에서 전체 루프가 검증된 것입니다.""",
    # 33 divider Part5
    """마지막으로 핵심 설계 결정을 정리하고 마무리합니다.""",
    # 34 결정 요약
    """이 도구를 만들며 내린 핵심 결정들입니다. stdout을 프로토콜 전용으로 보호한 것, 새 콘솔로 PTY를 우회한 것, rsync에 상대경로를 쓴 것, NativeSearch를 일시 적용한 것, PDK와 컴파일러를 자동 감지한 것, 32비트로 고정한 것, 그리고 PDK 판정에서 라이선스 요구를 뺀 것입니다. 관통하는 주제는 세 가지입니다 — 헤드리스 환경에서 동작할 것, 어느 PC에서나 이식될 것, 안전할 것.""",
    # 35 마무리
    """더 깊은 내용은 저장소 문서를 참고하세요. 사용자 매뉴얼, 설치·문제해결 가이드, MCP 내부 동작 설명, Skill 지식 구조가 각각 정리돼 있고 공개 저장소 주소도 있습니다. 기억할 한 가지는 — 그냥 물어보면 됩니다. 자연어로 작성부터 빌드, 다운로드, 제어까지. 이상으로 발표를 마칩니다.""",
]


# ================================================================ 조립
def build():
    s_cover()
    s_toc()
    s_divider("Part 1", "개요", "무엇을, 왜 — 문제의식과 전체 그림")
    s_problem(); s_solution(); s_two_axes(); s_sysdiagram()
    s_divider("Part 2", "Skill 계층", "Power PMAC을 '이해'하게 하는 지식")
    s_skill_mech(); s_skill_struct(); s_routing(); s_mental(); s_gotchas(); s_skill_use()
    s_divider("Part 3", "MCP 계층", "컨트롤러를 '조작'한다 — 동작 원리 핵심")
    s_mcp_what(); s_arch(); s_proto(); s_sequence(); s_tools(); s_bridge()
    s_build(); s_download(); s_live(); s_runtime(); s_pdkdetect(); s_x86()
    s_divider("Part 4", "설치 · 배포 · 운영", "setup.ps1 · 듀얼 리모트 · 워크플로 · 실증")
    s_setup(); s_deploy(); s_workflow(); s_proof()
    s_divider("Part 5", "마무리", "핵심 설계 결정 요약 · 참고자료")
    s_decisions(); s_closing()

    # 발표 노트(스크립트) 주입 — 슬라이드 순서 = 위 호출 순서
    for i, slide in enumerate(prs.slides):
        if i < len(NOTES) and NOTES[i]:
            set_notes(slide, NOTES[i])

    out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                       "PowerPMAC_MCP_교육자료.pptx")
    prs.save(out)
    print("OK: saved %d slides -> %s" % (len(prs.slides._sldIdLst), out))


if __name__ == "__main__":
    build()
